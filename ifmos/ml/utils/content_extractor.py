"""
Content Extraction Module
Handles PDF, Office documents, and images with OCR fallback
"""

import logging
from typing import Dict, List, Optional
from pathlib import Path
import mimetypes

try:
    import PyPDF2
    from pdf2image import convert_from_path
    import pdfplumber
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class ContentExtractor:
    """
    Universal content extractor for documents, PDFs, images.
    Integrates with OCR for scanned content.
    """

    def __init__(self, ocr_engine=None):
        """
        Initialize content extractor.

        Args:
            ocr_engine: Optional GPUOCREngine instance for scanned documents
        """
        self.logger = logging.getLogger(__name__)
        self.ocr_engine = ocr_engine

        # Verify dependencies
        if not PYPDF2_AVAILABLE:
            self.logger.warning("PyPDF2/pdfplumber not available. PDF support limited.")
        if not DOCX_AVAILABLE:
            self.logger.warning("python-docx not available. DOCX support disabled.")
        if not OPENPYXL_AVAILABLE:
            self.logger.warning("openpyxl not available. Excel support disabled.")

    def extract_content(self, file_path: str) -> Dict:
        """
        Extract content from any supported file type.

        Args:
            file_path: Path to file

        Returns:
            {
                'text': str,           # Extracted text
                'metadata': Dict,      # File metadata
                'method': str,         # Extraction method used
                'success': bool,       # Success status
                'page_count': int,     # Number of pages (if applicable)
                'confidence': float    # OCR confidence (if OCR used)
            }
        """
        file_path = Path(file_path)

        if not file_path.exists():
            return self._error_result(f"File not found: {file_path}")

        # Determine file type
        mime_type, _ = mimetypes.guess_type(str(file_path))
        extension = file_path.suffix.lower()

        self.logger.info(f"Extracting content from {file_path.name} ({mime_type})")

        try:
            # PDF files
            if extension == '.pdf':
                return self._extract_pdf(file_path)

            # Word documents
            elif extension in ['.docx', '.doc']:
                return self._extract_docx(file_path)

            # Excel files
            elif extension in ['.xlsx', '.xls']:
                return self._extract_excel(file_path)

            # Images
            elif extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif', '.webp']:
                return self._extract_image(file_path)

# Plain text and structured data            elif extension in [047.txt047, 047.log047, 047.csv047, 047.json047, 047.xml047, 047.md047, 047.rst047]:                return self._extract_text(file_path)            # Code files            elif extension in [047.py047, 047.js047, 047.ts047, 047.java047, 047.cpp047, 047.c047, 047.h047, 047.cs047, 047.go047, 047.rs047, 047.rb047, 047.php047]:                return self._extract_code(file_path)            # Script files            elif extension in [047.ps1047, 047.sh047, 047.bat047, 047.cmd047]:                return self._extract_script(file_path)            # Configuration files            elif extension in [047.yaml047, 047.yml047, 047.toml047, 047.ini047, 047.conf047, 047.config047]:                return self._extract_config(file_path)            # HTML files            elif extension in [047.html047, 047.htm047]:                return self._extract_html(file_path)            # PowerPoint            elif extension in [047.pptx047]:                return self._extract_powerpoint(file_path)
                return self._extract_text(file_path)

            else:
                return self._error_result(f"Unsupported file type: {extension}")

        except Exception as e:
            self.logger.error(f"Extraction failed for {file_path}: {e}")
            return self._error_result(str(e))

    def _extract_pdf(self, file_path: Path) -> Dict:
        """Extract text from PDF using PyPDF2 and pdfplumber, with OCR fallback."""
        if not PYPDF2_AVAILABLE:
            return self._error_result("PDF libraries not installed")

        text_content = []
        page_count = 0
        method = "digital_pdf"

        try:
            # Try digital PDF extraction with pdfplumber first (better for complex layouts)
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)

            # If minimal text extracted, likely scanned PDF - use OCR
            full_text = '\n'.join(text_content)
            if len(full_text.strip()) < 100 and self.ocr_engine:
                self.logger.info(f"Minimal text in PDF ({len(full_text)} chars), attempting OCR")
                return self._extract_pdf_with_ocr(file_path, page_count)

            return {
                'text': full_text,
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': 'pdf'
                },
                'method': method,
                'success': True,
                'page_count': page_count,
                'confidence': 1.0  # Digital PDF has perfect confidence
            }

        except Exception as e:
            self.logger.error(f"PDF extraction failed: {e}")
            # Try OCR as fallback
            if self.ocr_engine:
                return self._extract_pdf_with_ocr(file_path, 0)
            return self._error_result(str(e))

    def _extract_pdf_with_ocr(self, file_path: Path, page_count_hint: int = 0) -> Dict:
        """Extract text from scanned PDF using OCR."""
        if not self.ocr_engine:
            return self._error_result("OCR engine not available for scanned PDF")

        try:
            # Convert PDF pages to images
            images = convert_from_path(str(file_path), dpi=300)
            page_count = len(images)

            self.logger.info(f"OCR processing {page_count} pages from PDF")

            # Save images to temp files and run OCR
            import tempfile
            import os

            text_content = []
            total_confidence = 0

            with tempfile.TemporaryDirectory() as temp_dir:
                image_paths = []

                # Save all pages as images
                for i, image in enumerate(images):
                    img_path = os.path.join(temp_dir, f'page_{i+1}.png')
                    image.save(img_path, 'PNG')
                    image_paths.append(img_path)

                # Batch OCR processing for speed
                results = self.ocr_engine.extract_text_batch(image_paths)

                for result in results:
                    if result['success']:
                        text_content.append(result['text'])
                        total_confidence += result['confidence']

            avg_confidence = total_confidence / page_count if page_count > 0 else 0

            return {
                'text': '\n\n'.join(text_content),
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': 'pdf_scanned'
                },
                'method': 'ocr_pdf',
                'success': True,
                'page_count': page_count,
                'confidence': avg_confidence
            }

        except Exception as e:
            self.logger.error(f"PDF OCR extraction failed: {e}")
            return self._error_result(str(e))

    def _extract_docx(self, file_path: Path) -> Dict:
        """Extract text from Word documents."""
        if not DOCX_AVAILABLE:
            return self._error_result("python-docx not installed")

        try:
            doc = Document(file_path)
            text_content = []

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = '\t'.join(cell.text for cell in row.cells)
                    text_content.append(row_text)

            return {
                'text': '\n'.join(text_content),
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': 'docx',
                    'paragraph_count': len(doc.paragraphs),
                    'table_count': len(doc.tables)
                },
                'method': 'docx',
                'success': True,
                'page_count': 1,
                'confidence': 1.0
            }

        except Exception as e:
            self.logger.error(f"DOCX extraction failed: {e}")
            return self._error_result(str(e))

    def _extract_excel(self, file_path: Path) -> Dict:
        """Extract text from Excel files."""
        if not OPENPYXL_AVAILABLE:
            return self._error_result("openpyxl not installed")

        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== Sheet: {sheet_name} ===")

                for row in sheet.iter_rows(values_only=True):
                    row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_content.append(row_text)

            return {
                'text': '\n'.join(text_content),
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': 'excel',
                    'sheet_count': len(workbook.sheetnames),
                    'sheet_names': workbook.sheetnames
                },
                'method': 'excel',
                'success': True,
                'page_count': len(workbook.sheetnames),
                'confidence': 1.0
            }

        except Exception as e:
            self.logger.error(f"Excel extraction failed: {e}")
            return self._error_result(str(e))

    def _extract_image(self, file_path: Path) -> Dict:
        """Extract text from images using OCR."""
        if not self.ocr_engine:
            return self._error_result("OCR engine not available for images")

        try:
            result = self.ocr_engine.extract_text(str(file_path))

            if result['success']:
                return {
                    'text': result['text'],
                    'metadata': {
                        'file_name': file_path.name,
                        'file_size': file_path.stat().st_size,
                        'file_type': 'image'
                    },
                    'method': 'ocr_image',
                    'success': True,
                    'page_count': 1,
                    'confidence': result['confidence']
                }
            else:
                return self._error_result(result.get('error', 'OCR failed'))

        except Exception as e:
            self.logger.error(f"Image OCR failed: {e}")
            return self._error_result(str(e))

    def _extract_text(self, file_path: Path) -> Dict:
        """Extract content from plain text files."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            text_content = None

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue

            if text_content is None:
                return self._error_result("Failed to decode text file")

            return {
                'text': text_content,
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': 'text'
                },
                'method': 'text',
                'success': True,
                'page_count': 1,
                'confidence': 1.0
            }

        except Exception as e:
            self.logger.error(f"Text extraction failed: {e}")
            return self._error_result(str(e))

    def _error_result(self, error_msg: str) -> Dict:
        """Generate error result dictionary."""
        return {
            'text': '',
            'metadata': {},
            'method': 'error',
            'success': False,
            'page_count': 0,
            'confidence': 0.0,
            'error': error_msg
        }

    def extract_batch(self, file_paths: List[str]) -> List[Dict]:
        """
        Extract content from multiple files.

        Args:
            file_paths: List of file paths

        Returns:
            List of extraction results
        """
        results = []
        for file_path in file_paths:
            result = self.extract_content(file_path)
            results.append(result)

        return results


# Convenience function
def create_extractor(ocr_engine=None):
    """
    Factory function to create content extractor.

    Args:
        ocr_engine: Optional OCR engine instance

    Returns:
        Configured ContentExtractor
    """
    return ContentExtractor(ocr_engine=ocr_engine)

    def _extract_code(self, file_path: Path) -> Dict:
        """Extract content from source code files."""
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252']
            text_content = None
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            if text_content is None:
                return self._error_result('Failed to decode code file')
            lines = text_content.split('\n')
            return {
                'text': text_content,
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': f'code_{file_path.suffix[1:]}',
                    'line_count': len(lines)
                },
                'method': 'code',
                'success': True,
                'page_count': 1,
                'confidence': 1.0
            }
        except Exception as e:
            self.logger.error(f"Code extraction failed: {e}")
            return self._error_result(str(e))

    def _extract_script(self, file_path: Path) -> Dict:
        """Extract content from script files."""
        try:
            encodings = ['utf-8', 'utf-16', 'latin-1']
            text_content = None
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            if text_content is None:
                return self._error_result('Failed to decode script')
            return {
                'text': text_content,
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': f'script_{file_path.suffix[1:]}'
                },
                'method': 'script',
                'success': True,
                'page_count': 1,
                'confidence': 1.0
            }
        except Exception as e:
            self.logger.error(f"Script extraction failed: {e}")
            return self._error_result(str(e))

    def _extract_config(self, file_path: Path) -> Dict:
        """Extract content from config files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
            return {
                'text': text_content,
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': f'config_{file_path.suffix[1:]}'
                },
                'method': 'config',
                'success': True,
                'page_count': 1,
                'confidence': 1.0
            }
        except Exception as e:
            self.logger.error(f"Config extraction failed: {e}")
            return self._error_result(str(e))

    def _extract_html(self, file_path: Path) -> Dict:
        """Extract text from HTML files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                for script in soup(['script', 'style']):
                    script.decompose()
                text_content = soup.get_text(separator='\n', strip=True)
            except ImportError:
                import re
                text_content = re.sub('<script.*?</script>', '', html_content, flags=re.DOTALL)
                text_content = re.sub('<.*?>', '', text_content)
                
            return {
                'text': text_content,
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': 'html'
                },
                'method': 'html',
                'success': True,
                'page_count': 1,
                'confidence': 1.0
            }
        except Exception as e:
            self.logger.error(f"HTML extraction failed: {e}")
            return self._error_result(str(e))

    def _extract_powerpoint(self, file_path: Path) -> Dict:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        text_content.append(shape.text)
            return {
                'text': '\n'.join(text_content),
                'metadata': {
                    'file_name': file_path.name,
                    'file_size': file_path.stat().st_size,
                    'file_type': 'powerpoint',
                    'slide_count': len(prs.slides)
                },
                'method': 'powerpoint',
                'success': True,
                'page_count': len(prs.slides),
                'confidence': 1.0
            }
        except ImportError:
            return self._error_result('python-pptx not installed')
        except Exception as e:
            self.logger.error(f"PowerPoint extraction failed: {e}")
            return self._error_result(str(e))
